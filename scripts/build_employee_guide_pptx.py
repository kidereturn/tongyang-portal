# -*- coding: utf-8 -*-
"""
임직원 사용법 PPTX 생성 — 사이트 디자인 톤 (토스 브랜드 + 웜 톤)
docs/임직원_사용법_대본_슬라이드.md 의 내용을 기반으로 14장 슬라이드 생성.
실행: python scripts/build_employee_guide_pptx.py
출력: docs/임직원_사용법_안내.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ────────────────────────────────────────────
# 디자인 토큰 (사이트와 동일 톤)
# ────────────────────────────────────────────
BRAND_900 = RGBColor(0x0F, 0x1E, 0x36)  # 짙은 네이비 (제목)
BRAND_800 = RGBColor(0x1A, 0x2E, 0x4D)
BRAND_700 = RGBColor(0x31, 0x82, 0xF6)  # 토스 블루 (포인트)
BRAND_600 = RGBColor(0x4F, 0x96, 0xF8)
BRAND_500 = RGBColor(0x80, 0xB1, 0xFB)
BRAND_50  = RGBColor(0xEE, 0xF4, 0xFE)  # 연한 블루 카드 배경
WARM_900  = RGBColor(0x33, 0x33, 0x33)
WARM_700  = RGBColor(0x4E, 0x5A, 0x6F)  # 본문
WARM_500  = RGBColor(0x8B, 0x95, 0xA1)  # 보조 텍스트
WARM_400  = RGBColor(0xB0, 0xB8, 0xC1)
WARM_200  = RGBColor(0xE5, 0xE8, 0xEB)  # 라인
WARM_100  = RGBColor(0xF2, 0xF4, 0xF6)  # 카드 배경
WARM_50   = RGBColor(0xF9, 0xFA, 0xFB)  # 페이지 배경
AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)  # 강조 (PIN/경고)
AMBER_50  = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_600 = RGBColor(0x16, 0xA3, 0x4A)
RED_600   = RGBColor(0xDC, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 (10인치 x 5.625인치)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Pretendard"  # 동양/토스 톤. 미설치시 fallback
FONT_BODY  = "Pretendard"
FONT_MONO  = "JetBrains Mono"

# ────────────────────────────────────────────
# 헬퍼
# ────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *,
             font=FONT_BODY, size=14, bold=False, color=WARM_900,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_space=1.3):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_space
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box, tf

def add_card(slide, x, y, w, h, *, fill=WHITE, border=WARM_200, radius=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    if not shadow:
        sh.shadow.inherit = False
    # 코너 반경 (rounded rect adjustment)
    if radius and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = 0.08
        except Exception:
            pass
    return sh

def add_pill(slide, x, y, text, *, fill=BRAND_50, color=BRAND_700, size=10, bold=True, padding_x=0.18, padding_y=0.05):
    # 너비를 텍스트 길이에 비례하게 추정
    w = Inches(0.32 + 0.10 * len(text))
    h = Inches(0.32)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try: sh.adjustments[0] = 0.5
    except: pass
    tf = sh.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return sh, w

def add_eyebrow(slide, x, y, eyebrow, title):
    """ '카테고리 · 부제' + 큰 타이틀 (사이트 pg-head 톤) """
    add_text(slide, x, y, Inches(8), Inches(0.32),
             eyebrow, font=FONT_MONO, size=11, bold=True, color=BRAND_700)
    add_text(slide, x, y + Inches(0.32), Inches(11.5), Inches(0.9),
             title, font=FONT_TITLE, size=32, bold=True, color=BRAND_900, line_space=1.1)

def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "내부회계 포털 (TYIA) · 임직원 사용법", font=FONT_BODY, size=9, color=WARM_400)
    add_text(slide, Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{page_num} / {total}", font=FONT_MONO, size=9, color=WARM_400, align=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    # 좌측 브랜드 마크 + 타이틀 — 모든 본문 슬라이드 상단
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()
    bar.shadow.inherit = False
    # 좌측 브랜드 dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.18), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = BRAND_700
    dot.line.fill.background(); dot.shadow.inherit = False
    add_text(slide, Inches(0.85), Inches(0.13), Inches(6), Inches(0.3),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=10, bold=True, color=BRAND_900)
    # 하단 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.55), SLIDE_W, Emu(9525))  # 1pt
    line.fill.solid(); line.fill.fore_color.rgb = WARM_200
    line.line.fill.background(); line.shadow.inherit = False

# ────────────────────────────────────────────
# 슬라이드 빌더
# ────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

TOTAL_PAGES = 14

# ── 1. 표지 ───────────────────────────────────
def slide_cover():
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BRAND_900)
    # 좌측 큰 dot
    big_dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.7), Inches(0.42), Inches(0.42))
    big_dot.fill.solid(); big_dot.fill.fore_color.rgb = BRAND_700
    big_dot.line.fill.background(); big_dot.shadow.inherit = False
    add_text(s, Inches(1.25), Inches(0.78), Inches(8), Inches(0.4),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=12, bold=True, color=WHITE)

    # 메인 카피
    add_text(s, Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.5),
             "TYIA · 내부회계 포털", font=FONT_MONO, size=14, bold=True, color=BRAND_500)
    add_text(s, Inches(0.7), Inches(2.85), Inches(11.5), Inches(1.6),
             "임직원 사용법", font=FONT_TITLE, size=72, bold=True, color=WHITE, line_space=1.0)
    add_text(s, Inches(0.7), Inches(4.45), Inches(11.5), Inches(0.7),
             "5분이면 충분합니다.", font=FONT_TITLE, size=28, bold=False, color=BRAND_500, line_space=1.0)

    # 메타
    meta_y = Inches(6.3)
    add_text(s, Inches(0.7), meta_y, Inches(5), Inches(0.3),
             "(주)동양 내부회계관리실", font=FONT_BODY, size=12, color=WARM_400)
    add_text(s, Inches(0.7), meta_y + Inches(0.3), Inches(5), Inches(0.3),
             "발표일 · 발표자 정보", font=FONT_MONO, size=10, color=WARM_500)

    # URL 우측 하단
    add_text(s, Inches(8.5), Inches(6.6), Inches(4.2), Inches(0.4),
             "https://tyia.vercel.app", font=FONT_MONO, size=14, bold=True, color=BRAND_500, align=PP_ALIGN.RIGHT)

# ── 2. 한 줄 요약 ────────────────────────────
def slide_summary():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "OVERVIEW · 핵심 한 줄",
                "증빙 첨부, 결재상신, 학습 — 모두 한 곳에서.")

    # 3-step 카드
    step_y = Inches(2.8)
    step_w = Inches(3.95); step_h = Inches(2.7)
    gap = Inches(0.15)
    steps = [
        ("01", "포털 접속", "매일 1번", "본인 통제활동 확인"),
        ("02", "드래그 앤 드롭", "증빙 업로드", "PDF·Excel 자동 저장"),
        ("03", "결재상신", "한 번 클릭", "승인자에게 자동 전달"),
    ]
    for i, (n, title, sub, desc) in enumerate(steps):
        x = Inches(0.6) + (step_w + gap) * i
        add_card(s, x, step_y, step_w, step_h, fill=WHITE, border=WARM_200)
        add_text(s, x + Inches(0.4), step_y + Inches(0.3), Inches(2), Inches(0.5),
                 n, font=FONT_MONO, size=42, bold=True, color=BRAND_700, line_space=1.0)
        add_text(s, x + Inches(0.4), step_y + Inches(1.05), step_w - Inches(0.8), Inches(0.6),
                 title, font=FONT_TITLE, size=22, bold=True, color=BRAND_900)
        add_text(s, x + Inches(0.4), step_y + Inches(1.7), step_w - Inches(0.8), Inches(0.4),
                 sub, font=FONT_BODY, size=12, color=WARM_500)
        add_text(s, x + Inches(0.4), step_y + Inches(2.05), step_w - Inches(0.8), Inches(0.5),
                 desc, font=FONT_BODY, size=13, color=WARM_700)
    # 대본 인용
    add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.6),
             "“복잡한 시스템 아닙니다. 핵심은 세 단계 — 접속, 업로드, 상신. 이게 다입니다.”",
             font=FONT_TITLE, size=14, color=WARM_500, line_space=1.4)
    add_footer(s, 2, TOTAL_PAGES)

# ── 3. 로그인 ────────────────────────────────
def slide_login():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "STEP 01 · 로그인",
                "사번 + ty + 사번 — 그게 전부.")

    # 좌측 정보 카드
    card_x = Inches(0.6); card_y = Inches(2.7); card_w = Inches(7.5); card_h = Inches(4)
    add_card(s, card_x, card_y, card_w, card_h, fill=WHITE)
    add_text(s, card_x + Inches(0.4), card_y + Inches(0.4), card_w - Inches(0.8), Inches(0.4),
             "접속 정보", font=FONT_BODY, size=12, bold=True, color=WARM_500)

    rows = [
        ("URL",        "https://tyia.vercel.app", FONT_MONO),
        ("사번",        "6자리 사번 (예: 101267)", FONT_BODY),
        ("초기 비밀번호","ty + 사번 (예: ty101267)", FONT_MONO),
        ("변경 위치",   "마이페이지 → 비밀번호 변경", FONT_BODY),
    ]
    yy = card_y + Inches(0.95)
    for label, value, fnt in rows:
        add_text(s, card_x + Inches(0.4), yy, Inches(2.0), Inches(0.4),
                 label, font=FONT_BODY, size=12, color=WARM_500)
        add_text(s, card_x + Inches(2.5), yy, card_w - Inches(2.9), Inches(0.4),
                 value, font=fnt, size=15, bold=True, color=BRAND_900)
        yy += Inches(0.65)

    # 우측 강조 박스
    rx = Inches(8.4); ry = card_y; rw = Inches(4.3); rh = card_h
    add_card(s, rx, ry, rw, rh, fill=BRAND_50, border=BRAND_500)
    add_text(s, rx + Inches(0.4), ry + Inches(0.5), rw - Inches(0.8), Inches(0.5),
             "🔐  보안 안내", font=FONT_TITLE, size=16, bold=True, color=BRAND_700)
    add_text(s, rx + Inches(0.4), ry + Inches(1.2), rw - Inches(0.8), rh - Inches(1.5),
             "첫 로그인 후 반드시 비밀번호를\n변경해 주세요.\n\n· 8자 이상\n· 영문 + 숫자 혼합\n· 공용 PC 사용 시 로그아웃 필수",
             font=FONT_BODY, size=13, color=BRAND_900, line_space=1.6)

    add_footer(s, 3, TOTAL_PAGES)

# ── 4. 역할 안내 ─────────────────────────────
def slide_roles():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "ROLES · 역할 구조",
                "OWNER · CONTROLLER · ADMIN")

    role_y = Inches(2.8); role_h = Inches(3.6); gap = Inches(0.18)
    role_w = Inches(4.05)
    roles = [
        ("OWNER",      "담당자",   "442",  "본인 통제활동 증빙\n업로드 · 결재상신",         BRAND_700, BRAND_50),
        ("CONTROLLER", "승인자",   "33",   "담당자 상신 검토\n승인 · 반려 처리",            GREEN_600, RGBColor(0xE5,0xF7,0xEC)),
        ("ADMIN",      "관리자",   "3",    "전사 관리 · 모니터링\n메모 작성 (내부회계팀)",   AMBER_500, AMBER_50),
    ]
    for i, (key, label, count, desc, color, fill) in enumerate(roles):
        x = Inches(0.6) + (role_w + gap) * i
        add_card(s, x, role_y, role_w, role_h, fill=WHITE)
        # 좌측 컬러 바
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, role_y, Inches(0.12), role_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background(); bar.shadow.inherit = False
        # role key (mono)
        add_text(s, x + Inches(0.45), role_y + Inches(0.4), role_w - Inches(0.6), Inches(0.4),
                 key, font=FONT_MONO, size=11, bold=True, color=color)
        add_text(s, x + Inches(0.45), role_y + Inches(0.85), role_w - Inches(0.6), Inches(0.55),
                 label, font=FONT_TITLE, size=28, bold=True, color=BRAND_900)
        # 인원수
        add_text(s, x + Inches(0.45), role_y + Inches(1.65), role_w - Inches(0.6), Inches(0.4),
                 f"{count}명", font=FONT_MONO, size=20, bold=True, color=WARM_500)
        # 설명
        add_text(s, x + Inches(0.45), role_y + Inches(2.4), role_w - Inches(0.7), Inches(1.1),
                 desc, font=FONT_BODY, size=13, color=WARM_700, line_space=1.5)

    add_text(s, Inches(0.6), Inches(6.6), Inches(12.2), Inches(0.4),
             "본인 역할은 로그인 후 우측 상단에서 확인할 수 있습니다.",
             font=FONT_BODY, size=12, color=WARM_500)
    add_footer(s, 4, TOTAL_PAGES)

# ── 5. 담당자 사용법 (가장 중요) ─────────────
def slide_owner_flow():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "STEP 02 · 담당자 (Owner)",
                "증빙 업로드부터 결재상신까지.")

    # 7-step 좌측 리스트 카드
    list_x = Inches(0.6); list_y = Inches(2.8); list_w = Inches(8.0); list_h = Inches(4)
    add_card(s, list_x, list_y, list_w, list_h, fill=WHITE)
    steps = [
        ("증빙관리", "메뉴 클릭"),
        ("통제활동 목록", "본인 담당 활동 표시"),
        ("업로드 버튼", "→ 모달 창 열림"),
        ("드래그 앤 드롭", "또는 클릭 (PDF/Excel)"),
        ("자동 중간저장", "파일 1개마다 자동 ✨"),
        ("결재상신 버튼", "모든 모집단 완료 후"),
        ("자동 알림", "승인자에게 즉시 전송"),
    ]
    yy = list_y + Inches(0.35)
    for i, (t, sub) in enumerate(steps):
        # 동그라미 번호
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, list_x + Inches(0.4), yy + Inches(0.06), Inches(0.36), Inches(0.36))
        num.fill.solid(); num.fill.fore_color.rgb = BRAND_700
        num.line.fill.background(); num.shadow.inherit = False
        ntf = num.text_frame
        ntf.margin_left=Emu(0); ntf.margin_right=Emu(0); ntf.margin_top=Emu(0); ntf.margin_bottom=Emu(0)
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
        nr = np.add_run(); nr.text = str(i+1); nr.font.name = FONT_MONO; nr.font.size = Pt(11); nr.font.bold = True; nr.font.color.rgb = WHITE

        add_text(s, list_x + Inches(0.95), yy, Inches(3.2), Inches(0.4),
                 t, font=FONT_TITLE, size=14, bold=True, color=BRAND_900)
        add_text(s, list_x + Inches(4.2), yy, Inches(3.7), Inches(0.4),
                 sub, font=FONT_BODY, size=12, color=WARM_700)
        yy += Inches(0.50)

    # 우측 경고 카드
    wx = Inches(8.85); wy = list_y; ww = Inches(3.85); wh = list_h
    add_card(s, wx, wy, ww, wh, fill=AMBER_50, border=AMBER_500)
    add_text(s, wx + Inches(0.35), wy + Inches(0.4), ww - Inches(0.6), Inches(0.4),
             "⚠️  결재상신 조건", font=FONT_TITLE, size=15, bold=True, color=AMBER_500)
    add_text(s, wx + Inches(0.35), wy + Inches(1.1), ww - Inches(0.6), wh - Inches(1.4),
             "결재상신은 모든 모집단\n항목이 업로드된 뒤에만\n가능합니다.\n\n진행률(예: 8/10)이\n100% 가 되어야\n버튼이 활성화됩니다.",
             font=FONT_BODY, size=13, color=WARM_900, line_space=1.55)

    add_footer(s, 5, TOTAL_PAGES)

# ── 6. 담당자 FAQ ────────────────────────────
def slide_owner_faq():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "FAQ · 담당자",
                "자주 묻는 질문 4가지.")

    faqs = [
        ("Q1. 잘못 올렸어요. 바꿀 수 있나요?",
         "교체 버튼 → 새 파일 선택 → 중간저장 누르면 적용",
         BRAND_700),
        ("Q2. 삭제하고 다시 올리고 싶어요.",
         "삭제 버튼 → 즉시 DB·Storage 에서 제거 → 새로 업로드",
         BRAND_700),
        ("Q3. 상신 후에도 수정할 수 있나요?",
         "승인 전: 승인자가 반려해야 수정 가능\n승인 후: 관리자에게 ‘수정제출’ 요청",
         AMBER_500),
        ("Q4. 다른 사람 증빙 볼 수 있나요?",
         "❌ 절대 불가능. 본인 담당 활동만 표시됩니다.",
         RED_600),
    ]
    grid_x = Inches(0.6); grid_y = Inches(2.7); col_w = Inches(6.05); row_h = Inches(2.05); gap_x = Inches(0.18); gap_y = Inches(0.18)
    for i, (q, a, color) in enumerate(faqs):
        col = i % 2; row = i // 2
        x = grid_x + (col_w + gap_x) * col
        y = grid_y + (row_h + gap_y) * row
        add_card(s, x, y, col_w, row_h, fill=WHITE)
        # 질문
        add_text(s, x + Inches(0.4), y + Inches(0.32), col_w - Inches(0.7), Inches(0.45),
                 q, font=FONT_TITLE, size=15, bold=True, color=color)
        # 답
        add_text(s, x + Inches(0.4), y + Inches(0.95), col_w - Inches(0.7), row_h - Inches(1.1),
                 a, font=FONT_BODY, size=13, color=WARM_700, line_space=1.45)

    add_footer(s, 6, TOTAL_PAGES)

# ── 7. 승인자 사용법 ─────────────────────────
def slide_controller():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "STEP 03 · 승인자 (Controller)",
                "한 번의 클릭으로 검토 완료.")

    items = [
        ("01", "증빙관리", "본인 담당 통제활동 표시"),
        ("02", "상신 건만 노출", "담당자가 상신한 활동에만 액션 버튼"),
        ("03", "승인 / 반려", "사유 입력 시 담당자에게 자동 전달"),
        ("04", "다운로드", "개별 또는 전체 ZIP 다운로드 지원"),
    ]
    y = Inches(2.8); h = Inches(0.95); gap = Inches(0.16)
    for i, (n, title, desc) in enumerate(items):
        yy = y + (h + gap) * i
        add_card(s, Inches(0.6), yy, Inches(12.1), h, fill=WHITE)
        add_text(s, Inches(0.95), yy + Inches(0.22), Inches(0.9), Inches(0.55),
                 n, font=FONT_MONO, size=22, bold=True, color=BRAND_700)
        add_text(s, Inches(2.0), yy + Inches(0.18), Inches(3.5), Inches(0.5),
                 title, font=FONT_TITLE, size=18, bold=True, color=BRAND_900)
        add_text(s, Inches(5.7), yy + Inches(0.22), Inches(7.0), Inches(0.5),
                 desc, font=FONT_BODY, size=13, color=WARM_700)

    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             "“반려하실 때는 사유를 적어주셔야 담당자가 무엇을 수정할지 알 수 있어요.”",
             font=FONT_TITLE, size=12, color=WARM_500)
    add_footer(s, 7, TOTAL_PAGES)

# ── 8. 관리자 사용법 ─────────────────────────
def slide_admin():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "STEP 04 · 관리자 (Admin)",
                "검토결과 · 메모 · 수정제출.")

    # 2-column 카드
    left_x = Inches(0.6); left_y = Inches(2.8); left_w = Inches(6.05); col_h = Inches(4.1)
    right_x = Inches(6.85); right_w = Inches(5.85)

    add_card(s, left_x, left_y, left_w, col_h, fill=WHITE)
    add_text(s, left_x + Inches(0.4), left_y + Inches(0.4), left_w - Inches(0.8), Inches(0.4),
             "증빙관리 페이지", font=FONT_BODY, size=11, bold=True, color=WARM_500)
    add_text(s, left_x + Inches(0.4), left_y + Inches(0.85), left_w - Inches(0.8), Inches(0.5),
             "검토결과 + 메모", font=FONT_TITLE, size=22, bold=True, color=BRAND_900)
    items_left = [
        ("드롭다운", "미검토 / 검토중 / 완료 / 수정제출"),
        ("메모", "담당자·승인자에게 표시"),
        ("수정제출", "모든 승인 자동 취소 → 담당자 재상신"),
    ]
    yy = left_y + Inches(1.7)
    for k, v in items_left:
        add_text(s, left_x + Inches(0.4), yy, Inches(1.5), Inches(0.4),
                 k, font=FONT_BODY, size=12, bold=True, color=BRAND_700)
        add_text(s, left_x + Inches(1.95), yy, left_w - Inches(2.3), Inches(0.6),
                 v, font=FONT_BODY, size=13, color=WARM_700)
        yy += Inches(0.7)

    add_card(s, right_x, left_y, right_w, col_h, fill=BRAND_900)
    add_text(s, right_x + Inches(0.4), left_y + Inches(0.4), right_w - Inches(0.8), Inches(0.4),
             "관리자 페이지", font=FONT_MONO, size=11, bold=True, color=BRAND_500)
    add_text(s, right_x + Inches(0.4), left_y + Inches(0.85), right_w - Inches(0.8), Inches(0.5),
             "/admin", font=FONT_MONO, size=22, bold=True, color=WHITE)
    add_text(s, right_x + Inches(0.4), left_y + Inches(1.7), right_w - Inches(0.8), col_h - Inches(2.0),
             "· 사용자 마스터 관리\n· 통제활동 / 모집단 업로드\n· 진행률 모니터링\n· 공지사항 / 매뉴얼 발행\n· 강좌·웹툰·빙고·포인트 관리\n· 알림 발송 / AI 챗봇 관리",
             font=FONT_BODY, size=13, color=WARM_400, line_space=1.7)

    add_footer(s, 8, TOTAL_PAGES)

# ── 9. 부가 기능 ─────────────────────────────
def slide_extras():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "BEYOND CONTROLS · 부가 기능",
                "내부회계 외에도 학습과 소통을 한 곳에.")

    extras = [
        ("🎯", "빙고게임",   "매일 25문제 · 줄 완성 시 포인트"),
        ("🎓", "강좌",       "내부회계 교육 영상 학습"),
        ("📰", "회사소식",   "DART 공시 + 뉴스 자동 수집"),
        ("🤖", "AI 챗봇",    "내부회계 관련 질의응답 24시간"),
        ("💬", "Tell me",    "회사에 익명으로 제안하기"),
        ("🏅", "포인트",     "활동별 포인트 적립 + 월간 랭킹"),
    ]
    grid_x = Inches(0.6); grid_y = Inches(2.7); col_w = Inches(4.0); row_h = Inches(2.0)
    gap_x = Inches(0.18); gap_y = Inches(0.18)
    for i, (icon, title, desc) in enumerate(extras):
        col = i % 3; row = i // 3
        x = grid_x + (col_w + gap_x) * col
        y = grid_y + (row_h + gap_y) * row
        add_card(s, x, y, col_w, row_h, fill=WHITE)
        add_text(s, x + Inches(0.4), y + Inches(0.3), Inches(1), Inches(0.6),
                 icon, font=FONT_TITLE, size=28, color=BRAND_700)
        add_text(s, x + Inches(0.4), y + Inches(0.95), col_w - Inches(0.6), Inches(0.4),
                 title, font=FONT_TITLE, size=16, bold=True, color=BRAND_900)
        add_text(s, x + Inches(0.4), y + Inches(1.4), col_w - Inches(0.6), Inches(0.55),
                 desc, font=FONT_BODY, size=12, color=WARM_700, line_space=1.4)

    add_footer(s, 9, TOTAL_PAGES)

# ── 10. 도움이 필요할 때 ─────────────────────
def slide_help():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "SUPPORT · 도움이 필요할 때",
                "막히면 여기로.")

    helps = [
        ("🔔", "알림함",     "첫 화면 우측 상단 (벨 아이콘)",            BRAND_700),
        ("🤖", "AI 챗봇",    "/chatbot — 24시간 즉답",                   GREEN_600),
        ("💬", "Tell me",    "/tellme — 익명 문의",                      AMBER_500),
        ("☎️", "내부회계팀", "02 - 6150 - 7000 · 평일 09:00~18:00",     RED_600),
    ]
    y = Inches(2.8); h = Inches(0.95); gap = Inches(0.16)
    for i, (icon, title, desc, color) in enumerate(helps):
        yy = y + (h + gap) * i
        add_card(s, Inches(0.6), yy, Inches(12.1), h, fill=WHITE)
        # 좌측 컬러 dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), yy + Inches(0.32), Inches(0.32), Inches(0.32))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background(); dot.shadow.inherit = False
        add_text(s, Inches(1.55), yy + Inches(0.18), Inches(0.6), Inches(0.55),
                 icon, font=FONT_TITLE, size=20, color=color)
        add_text(s, Inches(2.3), yy + Inches(0.20), Inches(3.0), Inches(0.5),
                 title, font=FONT_TITLE, size=18, bold=True, color=BRAND_900)
        add_text(s, Inches(5.5), yy + Inches(0.22), Inches(7.0), Inches(0.5),
                 desc, font=FONT_BODY, size=13, color=WARM_700)

    add_footer(s, 10, TOTAL_PAGES)

# ── 11. 보안 안내 ────────────────────────────
def slide_security():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "SECURITY · 보안 안내",
                "안전한 사용을 위한 4가지 약속.")

    sec = [
        ("01", "비밀번호 정책", "8자 이상 · 영문 + 숫자 혼합 · 회사 공식 정책 준수"),
        ("02", "공용 PC 로그아웃", "회의실·교육장 PC 사용 시 반드시 로그아웃"),
        ("03", "피싱 방지",       "의심스러운 링크 클릭 금지 · IT팀에 즉시 신고"),
        ("04", "최소 권한",       "본인 데이터만 조회 가능 (RLS 시스템적 차단)"),
    ]
    y = Inches(2.8); h = Inches(0.97); gap = Inches(0.16)
    for i, (n, title, desc) in enumerate(sec):
        yy = y + (h + gap) * i
        add_card(s, Inches(0.6), yy, Inches(12.1), h, fill=WHITE)
        add_text(s, Inches(0.95), yy + Inches(0.22), Inches(0.9), Inches(0.55),
                 n, font=FONT_MONO, size=22, bold=True, color=BRAND_700)
        add_text(s, Inches(2.0), yy + Inches(0.20), Inches(3.6), Inches(0.5),
                 title, font=FONT_TITLE, size=17, bold=True, color=BRAND_900)
        add_text(s, Inches(5.7), yy + Inches(0.22), Inches(7.0), Inches(0.5),
                 desc, font=FONT_BODY, size=13, color=WARM_700)

    add_footer(s, 11, TOTAL_PAGES)

# ── 12. 시연 시나리오 ────────────────────────
def slide_demo():
    s = prs.slides.add_slide(blank); set_slide_bg(s, WARM_50); add_top_bar(s)
    add_eyebrow(s, Inches(0.6), Inches(0.85), "LIVE DEMO · 시연 시나리오",
                "5분 시연 — 9 단계.")

    demo = [
        ("로그인",          "30s"),
        ("홈 KPI / 검토 상태", "30s"),
        ("증빙관리 → 모달",  "30s"),
        ("PDF 드래그앤드롭",  "30s"),
        ("결재상신",        "30s"),
        ("관리자 검토 + 메모", "60s"),
        ("빙고 1게임",      "30s"),
        ("AI 챗봇 1회",     "30s"),
        ("로그아웃",        "15s"),
    ]
    grid_x = Inches(0.6); grid_y = Inches(2.7); col_w = Inches(4.0); row_h = Inches(1.25)
    gap_x = Inches(0.18); gap_y = Inches(0.16)
    for i, (title, dur) in enumerate(demo):
        col = i % 3; row = i // 3
        x = grid_x + (col_w + gap_x) * col
        y = grid_y + (row_h + gap_y) * row
        add_card(s, x, y, col_w, row_h, fill=WHITE)
        add_text(s, x + Inches(0.4), y + Inches(0.18), Inches(0.7), Inches(0.4),
                 f"{i+1:02d}", font=FONT_MONO, size=12, bold=True, color=BRAND_700)
        add_text(s, x + Inches(0.4), y + Inches(0.5), col_w - Inches(0.6), Inches(0.5),
                 title, font=FONT_TITLE, size=16, bold=True, color=BRAND_900)
        add_text(s, x + col_w - Inches(1.1), y + Inches(0.18), Inches(0.7), Inches(0.4),
                 dur, font=FONT_MONO, size=11, bold=True, color=WARM_500, align=PP_ALIGN.RIGHT)

    add_footer(s, 12, TOTAL_PAGES)

# ── 13. 한 장 요약 ─────────────────────────
def slide_summary_card():
    s = prs.slides.add_slide(blank); set_slide_bg(s, BRAND_900)
    add_text(s, Inches(0.7), Inches(0.7), Inches(8), Inches(0.4),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=11, bold=True, color=BRAND_500)

    add_text(s, Inches(0.7), Inches(1.3), Inches(11.5), Inches(1.0),
             "한 장으로 끝.", font=FONT_TITLE, size=44, bold=True, color=WHITE)

    # 4-column quick card
    items = [
        ("URL",      "tyia.vercel.app"),
        ("로그인",   "사번 / ty + 사번"),
        ("핵심",     "업로드 → 상신"),
        ("도움",     "AI 챗봇 / Tell me"),
    ]
    cy = Inches(3.4); cw = Inches(3.0); ch = Inches(2.6); gx = Inches(0.18)
    for i, (k, v) in enumerate(items):
        x = Inches(0.7) + (cw + gx) * i
        add_card(s, x, cy, cw, ch, fill=BRAND_800, border=None, radius=True)
        add_text(s, x + Inches(0.35), cy + Inches(0.35), cw - Inches(0.6), Inches(0.4),
                 k, font=FONT_MONO, size=11, bold=True, color=BRAND_500)
        add_text(s, x + Inches(0.35), cy + Inches(0.95), cw - Inches(0.6), ch - Inches(1.2),
                 v, font=FONT_TITLE, size=20, bold=True, color=WHITE, line_space=1.3)

    add_text(s, Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.5),
             "업무에 바로 쓰세요. 5분이면 충분합니다.",
             font=FONT_TITLE, size=18, color=BRAND_500)

# ── 14. Q&A ─────────────────────────────────
def slide_qa():
    s = prs.slides.add_slide(blank); set_slide_bg(s, BRAND_900)
    # 큰 Q&A
    add_text(s, Inches(0.7), Inches(2.6), Inches(11.5), Inches(2.0),
             "Q & A", font=FONT_TITLE, size=120, bold=True, color=WHITE, line_space=1.0,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.6), Inches(11.5), Inches(0.6),
             "사용 중 막힌 부분이 있다면 무엇이든 물어보세요.",
             font=FONT_TITLE, size=18, color=BRAND_500, align=PP_ALIGN.CENTER)
    # 하단 contact strip
    add_text(s, Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.4),
             "내부회계팀 · 02-6150-7000 · 평일 09:00 ~ 18:00",
             font=FONT_MONO, size=12, color=WARM_400, align=PP_ALIGN.CENTER)

# ────────────────────────────────────────────
# 빌드 실행
# ────────────────────────────────────────────
slide_cover()           # 1
slide_summary()         # 2
slide_login()           # 3
slide_roles()           # 4
slide_owner_flow()      # 5
slide_owner_faq()       # 6
slide_controller()      # 7
slide_admin()           # 8
slide_extras()          # 9
slide_help()            # 10
slide_security()        # 11
slide_demo()            # 12
slide_summary_card()    # 13
slide_qa()              # 14

OUT = os.path.join(os.path.dirname(__file__), '..', 'docs', '임직원_사용법_안내.pptx')
prs.save(OUT)
import sys
try:
    sys.stdout.reconfigure(encoding='utf-8')
except Exception:
    pass
out_abs = os.path.abspath(OUT)
print(f"[OK] saved: {out_abs}".encode('ascii', 'replace').decode('ascii'))
print(f"     slides={len(prs.slides)} 16:9 ({prs.slide_width.inches} x {prs.slide_height.inches} in)")
