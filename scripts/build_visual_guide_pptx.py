# -*- coding: utf-8 -*-
"""
사이트 캡쳐 기반 임직원 사용 가이드 PPTX
- 담당자(Owner) + 승인자(Controller) 워크플로우
- 각 단계: 스크린샷 + 클릭 영역 빨간 박스 + 단계 번호 라벨 + 설명 카드
- 출력: docs/임직원_사용가이드_시각.pptx
"""
import os, json, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

try:
    sys.stdout.reconfigure(encoding='utf-8')
except Exception:
    pass

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OWNER_DIR = os.path.join(ROOT, 'docs', 'captures', 'owner')
CTRL_DIR  = os.path.join(ROOT, 'docs', 'captures', 'controller')
ANNOT_DIR = os.path.join(ROOT, 'docs', 'captures', 'annotated')
os.makedirs(ANNOT_DIR, exist_ok=True)
OUT_PPTX  = os.path.join(ROOT, 'docs', '임직원_사용가이드_시각.pptx')

# 디자인 토큰 (사이트와 동일)
BRAND_900 = RGBColor(0x0F, 0x1E, 0x36)
BRAND_700 = RGBColor(0x31, 0x82, 0xF6)
BRAND_500 = RGBColor(0x80, 0xB1, 0xFB)
BRAND_50  = RGBColor(0xEE, 0xF4, 0xFE)
WARM_900  = RGBColor(0x33, 0x33, 0x33)
WARM_700  = RGBColor(0x4E, 0x5A, 0x6F)
WARM_500  = RGBColor(0x8B, 0x95, 0xA1)
WARM_400  = RGBColor(0xB0, 0xB8, 0xC1)
WARM_200  = RGBColor(0xE5, 0xE8, 0xEB)
WARM_100  = RGBColor(0xF2, 0xF4, 0xF6)
WARM_50   = RGBColor(0xF9, 0xFA, 0xFB)
RED_600   = RGBColor(0xDC, 0x26, 0x26)
GREEN_600 = RGBColor(0x16, 0xA3, 0x4A)
AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_50  = RGBColor(0xFE, 0xF3, 0xC7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333); SLIDE_H = Inches(7.5)
FONT_TITLE = "Pretendard"; FONT_BODY = "Pretendard"; FONT_MONO = "JetBrains Mono"

# ────────────────────────────────────────────
# PIL 이미지 마크업 — 빨간 박스 + 단계 번호
# ────────────────────────────────────────────
def get_font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/malgunbd.ttf" if bold else "C:/Windows/Fonts/malgun.ttf",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
    ]
    for fp in candidates:
        if os.path.exists(fp):
            try: return ImageFont.truetype(fp, size)
            except: pass
    return ImageFont.load_default()

def annotate(src_path, boxes, dst_path, *, label_size=44, box_color=(220, 38, 38), box_width=5):
    """
    boxes: [{x, y, w, h, label?, note?}] 단계 박스들
    빨간 라운드 박스 + 좌상단 원형 번호 + 옵셔널 노트
    """
    img = Image.open(src_path).convert('RGB')
    draw = ImageDraw.Draw(img, 'RGBA')
    for i, b in enumerate(boxes):
        x, y, w, h = b['x'], b['y'], b['w'], b['h']
        # 그림자 (살짝 어두운 외곽)
        for off in range(3, 0, -1):
            draw.rectangle([x-off, y-off, x+w+off, y+h+off],
                           outline=(220, 38, 38, 60), width=1)
        # 메인 박스 (빨간색 굵게)
        draw.rectangle([x, y, x+w, y+h], outline=box_color, width=box_width)

        # 단계 번호 원
        label = b.get('label', str(i+1))
        cx, cy = x - 14, y - 14
        r = 28
        draw.ellipse([cx-r, cy-r, cx+r, cy+r], fill=box_color, outline=(255,255,255), width=3)
        font = get_font(label_size if len(label) <= 2 else label_size-8, bold=True)
        # text 중앙 정렬
        bbox = draw.textbbox((0,0), label, font=font)
        tw = bbox[2]-bbox[0]; th = bbox[3]-bbox[1]
        draw.text((cx - tw/2, cy - th/2 - 4), label, fill=(255,255,255), font=font)

        # 옵셔널 노트 (박스 우상단)
        note = b.get('note')
        if note:
            font_note = get_font(22, bold=True)
            nb = draw.textbbox((0,0), note, font=font_note)
            nw = nb[2]-nb[0]; nh = nb[3]-nb[1]
            pad = 10
            nx = x + w + 12; ny = y - 4
            draw.rectangle([nx, ny, nx+nw+pad*2, ny+nh+pad*2], fill=(220,38,38), outline=None)
            draw.text((nx+pad, ny+pad-3), note, fill=(255,255,255), font=font_note)

    img.save(dst_path, 'PNG', optimize=True)
    return dst_path

def crop_around(src_path, box, dst_path, *, padding=80, draw_box=True):
    """ 클릭 영역 주위만 cropped 해서 확대용으로 저장 """
    img = Image.open(src_path).convert('RGB')
    W, H = img.size
    x1 = max(0, box['x'] - padding)
    y1 = max(0, box['y'] - padding)
    x2 = min(W, box['x'] + box['w'] + padding)
    y2 = min(H, box['y'] + box['h'] + padding)
    crop = img.crop((x1, y1, x2, y2))
    if draw_box:
        # 잘라낸 좌표계로 박스 다시 그리기
        d = ImageDraw.Draw(crop, 'RGBA')
        nx = box['x'] - x1; ny = box['y'] - y1
        for off in range(3, 0, -1):
            d.rectangle([nx-off, ny-off, nx+box['w']+off, ny+box['h']+off], outline=(220,38,38,80), width=1)
        d.rectangle([nx, ny, nx+box['w'], ny+box['h']], outline=(220,38,38), width=5)
    crop.save(dst_path, 'PNG', optimize=True)
    return dst_path

# ────────────────────────────────────────────
# PPTX 빌더 헬퍼
# ────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14, bold=False,
             color=WARM_900, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_space=1.3):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = line_space
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return box

def add_card(slide, x, y, w, h, *, fill=WHITE, border=WARM_200, radius=True):
    sh_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(sh_type, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(0.75)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    if radius:
        try: sh.adjustments[0] = 0.04
        except: pass
    return sh

def add_top_bar(slide, *, role_label=None, role_color=BRAND_700):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background(); bar.shadow.inherit = False
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(0.18), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = BRAND_700
    dot.line.fill.background(); dot.shadow.inherit = False
    add_text(slide, Inches(0.8), Inches(0.13), Inches(7), Inches(0.3),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=10, bold=True, color=BRAND_900)
    if role_label:
        # 우측 role pill
        pill_w = Inches(1.2); pill_h = Inches(0.35); pill_x = SLIDE_W - pill_w - Inches(0.6); pill_y = Inches(0.1)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_x, pill_y, pill_w, pill_h)
        pill.fill.solid(); pill.fill.fore_color.rgb = role_color
        pill.line.fill.background(); pill.shadow.inherit = False
        try: pill.adjustments[0] = 0.5
        except: pass
        tf = pill.text_frame; tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = role_label
        r.font.name = FONT_MONO; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.55), SLIDE_W, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = WARM_200
    line.line.fill.background(); line.shadow.inherit = False

def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "내부회계 포털 (TYIA) · 시각 가이드", font=FONT_BODY, size=9, color=WARM_400)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page_num} / {total}", font=FONT_MONO, size=9, color=WARM_400, align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, x, y, eyebrow, title, *, accent=BRAND_700):
    add_text(slide, x, y, Inches(8), Inches(0.32),
             eyebrow, font=FONT_MONO, size=11, bold=True, color=accent)
    add_text(slide, x, y + Inches(0.3), Inches(11.5), Inches(0.85),
             title, font=FONT_TITLE, size=28, bold=True, color=BRAND_900, line_space=1.1)

def add_step_slide(prs, *, page, total, eyebrow, title, image_path, captions, role_label=None, role_color=BRAND_700, accent=BRAND_700, image_max_w_in=8.4):
    """
    캡쳐 슬라이드: 좌측 2/3 큰 이미지 + 우측 1/3 단계 설명 카드
    captions: [(번호, 제목, 설명)]
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WARM_50); add_top_bar(s, role_label=role_label, role_color=role_color)
    add_eyebrow(s, Inches(0.6), Inches(0.85), eyebrow, title, accent=accent)

    # 이미지 — 비율 맞추기
    if os.path.exists(image_path):
        img = Image.open(image_path); iw, ih = img.size
        max_w_in = image_max_w_in
        max_h_in = 5.4
        ratio = iw / ih
        if max_w_in / ratio <= max_h_in:
            disp_w = Inches(max_w_in); disp_h = Inches(max_w_in / ratio)
        else:
            disp_h = Inches(max_h_in); disp_w = Inches(max_h_in * ratio)
        img_x = Inches(0.6); img_y = Inches(2.05)
        # 이미지 카드 배경
        add_card(s, img_x, img_y, disp_w, disp_h, fill=WHITE)
        s.shapes.add_picture(image_path, img_x, img_y, disp_w, disp_h)

    # 우측 단계 설명 카드
    card_x = Inches(9.2); card_y = Inches(2.05); card_w = Inches(3.55); card_h = Inches(5.0)
    add_card(s, card_x, card_y, card_w, card_h, fill=WHITE)
    add_text(s, card_x + Inches(0.35), card_y + Inches(0.4), card_w - Inches(0.7), Inches(0.4),
             "단계", font=FONT_BODY, size=11, bold=True, color=accent)

    yy = card_y + Inches(0.95)
    for n, t, desc in captions:
        # 번호 원
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, card_x + Inches(0.35), yy + Inches(0.02), Inches(0.36), Inches(0.36))
        num.fill.solid(); num.fill.fore_color.rgb = accent
        num.line.fill.background(); num.shadow.inherit = False
        ntf = num.text_frame
        for k in ['margin_left','margin_right','margin_top','margin_bottom']: setattr(ntf, k, Emu(0))
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
        nr = np.add_run(); nr.text = str(n)
        nr.font.name = FONT_MONO; nr.font.size = Pt(11); nr.font.bold = True; nr.font.color.rgb = WHITE

        add_text(s, card_x + Inches(0.85), yy, card_w - Inches(1.0), Inches(0.4),
                 t, font=FONT_TITLE, size=13, bold=True, color=BRAND_900)
        add_text(s, card_x + Inches(0.85), yy + Inches(0.4), card_w - Inches(1.05), Inches(1.5),
                 desc, font=FONT_BODY, size=11, color=WARM_700, line_space=1.45)
        yy += Inches(1.10)

    add_footer(s, page, total)

# ────────────────────────────────────────────
# 캡쳐 메타 + 박스 정의
# ────────────────────────────────────────────
def load_meta(p):
    if not os.path.exists(p): return {}
    with open(p, 'r', encoding='utf-8') as f: return json.load(f)

owner_01 = load_meta(os.path.join(OWNER_DIR, '01-meta.json'))
owner_03 = load_meta(os.path.join(OWNER_DIR, '03-meta.json'))
owner_05 = load_meta(os.path.join(OWNER_DIR, '05-meta.json'))
owner_07 = load_meta(os.path.join(OWNER_DIR, '07-meta.json'))
owner_09 = load_meta(os.path.join(OWNER_DIR, '09-meta.json'))

ctrl_04  = load_meta(os.path.join(CTRL_DIR, '04-meta.json'))
ctrl_06  = load_meta(os.path.join(CTRL_DIR, '06-meta.json'))
ctrl_07  = load_meta(os.path.join(CTRL_DIR, '07-meta.json'))

# 마크업 이미지 생성
def make_annotated(src, dst, boxes_with_labels):
    boxes = []
    for box, label, note in boxes_with_labels:
        if not box: continue
        b = dict(box); b['label'] = label
        if note: b['note'] = note
        boxes.append(b)
    if not boxes:
        # 박스 없으면 원본 복사
        Image.open(src).save(dst, 'PNG'); return dst
    return annotate(src, boxes, dst)

# Owner annotated
own_imgs = {}
src01 = os.path.join(OWNER_DIR, '01b-login-filled.png')
own_imgs['01'] = make_annotated(src01, os.path.join(ANNOT_DIR, 'own-01.png'),
    [(owner_01.get('idBox'), '1', '사번'), (owner_01.get('submitBox'), '2', '로그인')])

own_imgs['02'] = os.path.join(OWNER_DIR, '02-home.png')

src03 = os.path.join(OWNER_DIR, '03-home-with-gnb.png')
own_imgs['03'] = make_annotated(src03, os.path.join(ANNOT_DIR, 'own-03.png'),
    [(owner_03.get('evidenceMenuBox'), '1', '클릭')])

own_imgs['04'] = os.path.join(OWNER_DIR, '04-evidence-list.png')

src05 = os.path.join(OWNER_DIR, '05-evidence-confirm-btn.png')
own_imgs['05'] = make_annotated(src05, os.path.join(ANNOT_DIR, 'own-05.png'),
    [(owner_05.get('confirmBox'), '1', '클릭')])

own_imgs['06'] = os.path.join(OWNER_DIR, '06-modal-opened.png')

src07 = os.path.join(OWNER_DIR, '07-dropzone.png')
own_imgs['07'] = make_annotated(src07, os.path.join(ANNOT_DIR, 'own-07.png'),
    [(owner_07.get('dropZoneBox'), '1', '드래그앤드롭')])

own_imgs['08'] = os.path.join(OWNER_DIR, '08-files-existing.png') if os.path.exists(os.path.join(OWNER_DIR, '08-files-existing.png')) else os.path.join(OWNER_DIR, '08-uploaded.png')

src09 = os.path.join(OWNER_DIR, '09-submit-btn.png')
own_imgs['09'] = make_annotated(src09, os.path.join(ANNOT_DIR, 'own-09.png'),
    [(owner_09.get('submitBtnBox'), '1', '결재상신')])

# Controller annotated
ctrl_imgs = {}
ctrl_imgs['02'] = os.path.join(CTRL_DIR, '02-home-controller.png')
ctrl_imgs['04a'] = os.path.join(CTRL_DIR, '04-evidence-controller.png')
src_ctrl_04b = os.path.join(CTRL_DIR, '04b-approve-btn.png')
ctrl_imgs['04b'] = make_annotated(src_ctrl_04b, os.path.join(ANNOT_DIR, 'ctrl-04b.png'),
    [(ctrl_04.get('approveBtnBox'), '1', '승인 / 반려')])
ctrl_imgs['05'] = os.path.join(CTRL_DIR, '05-modal-controller.png')
src_ctrl_06 = os.path.join(CTRL_DIR, '06-download-btn.png')
ctrl_imgs['06'] = make_annotated(src_ctrl_06, os.path.join(ANNOT_DIR, 'ctrl-06.png'),
    [(ctrl_06.get('dlBtnBox'), '1', '다운로드')])
src_ctrl_07 = os.path.join(CTRL_DIR, '07-modal-footer.png')
ctrl_imgs['07'] = make_annotated(src_ctrl_07, os.path.join(ANNOT_DIR, 'ctrl-07.png'),
    [(ctrl_07.get('zipBox'), '1', '전체 ZIP'),
     (ctrl_07.get('approveModalBox'), '2', '승인'),
     (ctrl_07.get('rejectModalBox'), '3', '반려')])

# ────────────────────────────────────────────
# PPTX 빌드
# ────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# 표지 + 섹션 표지 + 9 owner + 6 controller + 마무리 = 약 18장
TOTAL = 18

# ── 표지 ─────────────────────────────────────
def cover():
    s = prs.slides.add_slide(blank); set_bg(s, BRAND_900)
    big = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.7), Inches(0.42), Inches(0.42))
    big.fill.solid(); big.fill.fore_color.rgb = BRAND_700; big.line.fill.background(); big.shadow.inherit = False
    add_text(s, Inches(1.25), Inches(0.78), Inches(8), Inches(0.4),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=12, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.5),
             "TYIA · 시각 사용 가이드", font=FONT_MONO, size=14, bold=True, color=BRAND_500)
    add_text(s, Inches(0.7), Inches(2.85), Inches(11.5), Inches(1.6),
             "보면 따라할 수 있는\n증빙 업무 매뉴얼", font=FONT_TITLE, size=58, bold=True, color=WHITE, line_space=1.05)
    add_text(s, Inches(0.7), Inches(5.5), Inches(11.5), Inches(0.5),
             "담당자 (OWNER) · 승인자 (CONTROLLER) 화면 캡쳐 기반",
             font=FONT_TITLE, size=18, color=BRAND_500)
    add_text(s, Inches(0.7), Inches(6.6), Inches(7), Inches(0.4),
             "(주)동양 내부회계관리실", font=FONT_BODY, size=12, color=WARM_400)
    add_text(s, Inches(8.5), Inches(6.6), Inches(4.2), Inches(0.4),
             "https://tyia.vercel.app", font=FONT_MONO, size=14, bold=True, color=BRAND_500, align=PP_ALIGN.RIGHT)

# ── 섹션 표지: OWNER ────────────────────────
def section_cover(title_eyebrow, title_main, label, color):
    s = prs.slides.add_slide(blank); set_bg(s, WARM_50); add_top_bar(s)
    # 큰 컬러 dot
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(2.4), Inches(0.6), Inches(0.6))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background(); d.shadow.inherit = False
    add_text(s, Inches(1.4), Inches(2.45), Inches(11), Inches(0.5),
             title_eyebrow, font=FONT_MONO, size=14, bold=True, color=color)
    add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5),
             title_main, font=FONT_TITLE, size=66, bold=True, color=BRAND_900, line_space=1.0)
    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
             label, font=FONT_TITLE, size=22, color=WARM_500)

# ── 캡쳐 슬라이드 (대량 함수) ────────────────
def slide_step(page, eyebrow, title, image, captions, role_label, role_color, accent=BRAND_700):
    add_step_slide(prs, page=page, total=TOTAL, eyebrow=eyebrow, title=title,
                   image_path=image, captions=captions,
                   role_label=role_label, role_color=role_color, accent=accent)

# 빌드
cover()                                                                         # 1
section_cover("PART 1 · 담당자", "OWNER 워크플로우", "OWNER · 442명",
              BRAND_700)                                                        # 2
slide_step(3, "OWNER · 01", "사번과 비밀번호로 로그인",
           own_imgs['01'],
           [(1, "사번 입력", "회사 사번 6자리 (예: 101267)"),
            (2, "비밀번호", "ty + 사번 (첫 로그인용)"),
            (3, "로그인 클릭", "엔터키도 가능")],
           "OWNER", BRAND_700)
slide_step(4, "OWNER · 02", "홈 화면 — 내 KPI · 검토 상태",
           own_imgs['02'],
           [(1, "내 점수", "이번 달 적립 포인트와 KPI 점수"),
            (2, "공지·매뉴얼", "관리자가 공유한 최신 정보"),
            (3, "GNB 메뉴", "상단 네비게이션 항상 노출")],
           "OWNER", BRAND_700)
slide_step(5, "OWNER · 03", "GNB 에서 ‘증빙관리’ 클릭",
           own_imgs['03'],
           [(1, "증빙관리 메뉴", "모든 업무는 여기서 시작"),
            (2, "내 담당만 표시", "다른 사람 활동은 보이지 않음")],
           "OWNER", BRAND_700)
slide_step(6, "OWNER · 04", "증빙목록 — 내 통제활동 한눈에",
           own_imgs['04'],
           [(1, "상태별 카드", "전체·상신완료·승인·수정제출·미완료"),
            (2, "필터 칩", "원하는 상태만 골라보기"),
            (3, "새로고침", "최신 상태로 즉시 동기화"),
            (4, "엑셀 다운로드", "전체 목록 한번에 저장")],
           "OWNER", BRAND_700)
slide_step(7, "OWNER · 05", "각 행의 ‘증빙확인’ 버튼 클릭",
           own_imgs['05'],
           [(1, "증빙확인", "통제활동별 상세 모달 열기"),
            (2, "건수 컬럼", "업로드된 모집단 / 전체 모집단")],
           "OWNER", BRAND_700)
slide_step(8, "OWNER · 06", "모달 열림 — 모집단 + 파일 영역",
           own_imgs['06'],
           [(1, "모집단 표", "각 행이 증빙을 첨부할 단위"),
            (2, "파일 박스", "업로드 결과 + 진행률 표시"),
            (3, "엑셀·ZIP 다운로드", "기존 자료 일괄 받기")],
           "OWNER", BRAND_700)
slide_step(9, "OWNER · 07", "드롭존에 파일 드래그앤드롭",
           own_imgs['07'],
           [(1, "드롭존", "여기에 파일을 끌어다 놓기"),
            (2, "PDF·Excel 지원", "다중 선택 가능"),
            (3, "자동 중간저장", "파일 추가 시 즉시 저장 ✨")],
           "OWNER", BRAND_700)
slide_step(10, "OWNER · 08", "업로드 후 — 다운로드·교체·삭제",
           own_imgs['08'],
           [(1, "다운로드", "본인 업로드 파일 다시 받기"),
            (2, "교체 / 삭제", "잘못 올렸을 때 즉시 수정"),
            (3, "진행률", "100% 도달 시 결재상신 활성화")],
           "OWNER", BRAND_700)
slide_step(11, "OWNER · 09", "결재상신 — 한 번 클릭으로 끝",
           own_imgs['09'],
           [(1, "결재상신 / 중간저장", "모든 모집단 완료 후 활성화"),
            (2, "자동 알림", "승인자에게 즉시 전달")],
           "OWNER", BRAND_700)

# Section: Controller
section_cover("PART 2 · 승인자", "CONTROLLER 워크플로우", "CONTROLLER · 33명",
              GREEN_600)                                                        # 12

slide_step(13, "CONTROLLER · 01", "홈 — CONTROLLER 라벨 확인",
           ctrl_imgs['02'],
           [(1, "역할 표시", "우측 상단에 CONTROLLER"),
            (2, "검토 대기 카드", "결재 대기 건수 한눈에")],
           "CONTROLLER", GREEN_600, accent=GREEN_600)
slide_step(14, "CONTROLLER · 02", "증빙관리 — 본인 담당 통제활동",
           ctrl_imgs['04a'],
           [(1, "내가 검토할 활동만", "담당자 상신완료된 건"),
            (2, "결재 컬럼", "행별 ‘승인 / 반려’ 노출")],
           "CONTROLLER", GREEN_600, accent=GREEN_600)
slide_step(15, "CONTROLLER · 03", "행 단위 ‘승인 / 반려’ 처리",
           ctrl_imgs['04b'],
           [(1, "승인", "즉시 ‘승인완료’ 로 전환"),
            (2, "반려", "사유 입력 → 담당자 알림"),
            (3, "재상신", "담당자가 수정 후 다시 상신")],
           "CONTROLLER", GREEN_600, accent=GREEN_600)
slide_step(16, "CONTROLLER · 04", "증빙확인 모달 — 파일 검토",
           ctrl_imgs['05'],
           [(1, "모집단 + 파일", "업로드된 증빙 일괄 확인"),
            (2, "다운로드", "개별 파일 받아 보기"),
            (3, "전체 ZIP", "모든 첨부 한번에 받기")],
           "CONTROLLER", GREEN_600, accent=GREEN_600)
slide_step(17, "CONTROLLER · 05", "다운로드 / ZIP / 모달 액션",
           ctrl_imgs['07'],
           [(1, "전체 ZIP 다운로드", "모든 증빙 한번에 저장"),
            (2, "승인", "모달 안에서도 즉시 처리"),
            (3, "반려", "사유 입력창 표시")],
           "CONTROLLER", GREEN_600, accent=GREEN_600)

# 마지막 — 5단계 요약
def slide_recap():
    s = prs.slides.add_slide(blank); set_bg(s, BRAND_900)
    add_text(s, Inches(0.7), Inches(0.7), Inches(8), Inches(0.4),
             "TONGYANG · INTERNAL CONTROLS", font=FONT_MONO, size=11, bold=True, color=BRAND_500)
    add_text(s, Inches(0.7), Inches(1.3), Inches(11.5), Inches(1.0),
             "이 5장만 외워도 충분합니다.", font=FONT_TITLE, size=42, bold=True, color=WHITE)
    items = [
        ("01", "tyia.vercel.app", "사번 / ty+사번"),
        ("02", "증빙관리 → 증빙확인", "행 클릭 → 모달 열기"),
        ("03", "드래그앤드롭", "파일 끌어 놓기"),
        ("04", "결재상신", "버튼 한 번"),
        ("05", "검토·승인", "승인자 영역"),
    ]
    cy = Inches(3.6); cw = Inches(2.4); ch = Inches(2.5); gx = Inches(0.12)
    for i, (n, t, sub) in enumerate(items):
        x = Inches(0.7) + (cw + gx) * i
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch)
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x4D)
        c.line.fill.background(); c.shadow.inherit = False
        try: c.adjustments[0] = 0.05
        except: pass
        add_text(s, x + Inches(0.3), cy + Inches(0.3), cw - Inches(0.6), Inches(0.5),
                 n, font=FONT_MONO, size=22, bold=True, color=BRAND_500)
        add_text(s, x + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(0.6),
                 t, font=FONT_TITLE, size=15, bold=True, color=WHITE, line_space=1.2)
        add_text(s, x + Inches(0.3), cy + Inches(1.75), cw - Inches(0.6), Inches(0.6),
                 sub, font=FONT_BODY, size=11, color=BRAND_500, line_space=1.3)
    add_text(s, Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.4),
             "막히면 — AI 챗봇 / Tell me / 02-6150-7000",
             font=FONT_MONO, size=12, color=WARM_400, align=PP_ALIGN.CENTER)

slide_recap()                                                                   # 18

prs.save(OUT_PPTX)
print(f"[OK] saved: {os.path.abspath(OUT_PPTX)}".encode('ascii', 'replace').decode('ascii'))
print(f"     slides={len(prs.slides)} 16:9")
